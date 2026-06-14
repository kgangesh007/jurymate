import os
import io
import re
import base64
import zipfile
import fitz                        # PyMuPDF
from pptx import Presentation
from PIL import Image, ImageDraw
from langchain_text_splitters import RecursiveCharacterTextSplitter, Language
import chromadb
from chromadb.utils.embedding_functions import DefaultEmbeddingFunction
import ollama
from core.database import (
    register_document, update_document_status, add_team
)

CHROMA_PATH     = os.path.join(os.path.dirname(__file__), "../data/chroma")
COLLECTION_NAME = "jurymate"

# ── ChromaDB setup ─────────────────────────────────────────
_chroma_client = None
_collection    = None

def get_collection():
    global _chroma_client, _collection
    if _collection is None:
        os.makedirs(CHROMA_PATH, exist_ok=True)
        _chroma_client = chromadb.PersistentClient(path=CHROMA_PATH)
        ef = DefaultEmbeddingFunction()
        _collection = _chroma_client.get_or_create_collection(
            name=COLLECTION_NAME,
            embedding_function=ef,
            metadata={"hnsw:space": "cosine"}
        )
    return _collection

# ── Vision helpers ─────────────────────────────────────────
def _is_vision_available() -> bool:
    try:
        models = ollama.list()
        return any(
            "llama3.2-vision" in m["name"].lower()
            for m in models.get("models", [])
        )
    except Exception:
        return False

def _image_to_base64(img: Image.Image) -> str:
    buffer = io.BytesIO()
    img.save(buffer, format="PNG")
    return base64.b64encode(buffer.getvalue()).decode("utf-8")

def _describe_slide_image(img: Image.Image, slide_num: int) -> str:
    try:
        img_b64 = _image_to_base64(img)
        resp = ollama.chat(
            model="llama3.2-vision:11b",
            messages=[{
                "role": "user",
                "content": f"""This is slide {slide_num} from a hackathon project presentation.
Describe everything visible in detail:
- All text on the slide
- Architecture diagrams and their components
- Technology names, tools, frameworks
- Flow diagrams and what they represent
- Any metrics, accuracy numbers, results
- Dataset names or model names
Be thorough and specific.""",
                "images": [img_b64]
            }]
        )
        return resp["message"]["content"]
    except Exception as e:
        return f"[Vision analysis unavailable: {str(e)}]"

def _slide_to_pil_image(slide, width=1280, height=720) -> Image.Image:
    img  = Image.new("RGB", (width, height), color=(255, 255, 255))
    draw = ImageDraw.Draw(img)
    y_pos = 30
    for shape in slide.shapes:
        if hasattr(shape, "text_frame"):
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    words, line = text.split(), ""
                    for word in words:
                        if len(line + word) < 90:
                            line += word + " "
                        else:
                            draw.text((30, y_pos), line.strip(), fill=(0,0,0))
                            y_pos += 22
                            line = word + " "
                    if line.strip():
                        draw.text((30, y_pos), line.strip(), fill=(0,0,0))
                        y_pos += 22
                    y_pos += 8
        if shape.shape_type == 13:
            try:
                pic = Image.open(io.BytesIO(shape.image.blob))
                pic.thumbnail((width-60, max(100, height-y_pos-30)), Image.LANCZOS)
                img.paste(pic, (30, y_pos))
                y_pos += pic.height + 15
            except Exception:
                pass
        if y_pos > height - 30:
            break
    return img

def _slide_has_images(slide) -> bool:
    return any(shape.shape_type == 13 for shape in slide.shapes)

def _slide_text_is_thin(slide, threshold=80) -> bool:
    text = " ".join(
        shape.text for shape in slide.shapes
        if hasattr(shape, "text")
    ).strip()
    return len(text) < threshold

# ── Code helpers ───────────────────────────────────────────
# Supported code file extensions
CODE_EXTENSIONS = {
    ".py":   "python",
    ".js":   "javascript",
    ".ts":   "typescript",
    ".java": "java",
    ".go":   "go",
    ".rs":   "rust",
    ".cpp":  "cpp",
    ".c":    "c",
    ".sh":   "bash",
    ".sql":  "sql",
    ".yaml": "yaml",
    ".yml":  "yaml",
    ".json": "json",
    ".md":   "markdown",
    ".txt":  "text",
    ".html": "html",
    ".css":  "css",
}

SKIP_DIRS = {
    "__pycache__", ".git", "node_modules", "venv", ".venv",
    "env", "dist", "build", ".idea", ".vscode"
}

def _get_file_language(filename: str) -> str:
    ext = os.path.splitext(filename)[1].lower()
    return CODE_EXTENSIONS.get(ext, "")

def _extract_code_structure(code: str, language: str, filename: str) -> str:
    """
    Extract meaningful structure from code:
    - imports/dependencies
    - function/class names
    - brief description of what each does
    """
    lines     = code.split("\n")
    imports   = []
    functions = []
    classes   = []

    for line in lines:
        stripped = line.strip()
        # Python imports
        if language == "python":
            if stripped.startswith(("import ", "from ")):
                imports.append(stripped)
            elif stripped.startswith("def "):
                func_name = stripped.split("(")[0].replace("def ", "")
                functions.append(func_name)
            elif stripped.startswith("class "):
                class_name = stripped.split("(")[0].split(":")[0].replace("class ", "")
                classes.append(class_name)
        # JavaScript/TypeScript
        elif language in ("javascript", "typescript"):
            if stripped.startswith(("import ", "require(")):
                imports.append(stripped[:100])
            elif "function " in stripped or "const " in stripped and "=>" in stripped:
                functions.append(stripped[:80])
        # Java
        elif language == "java":
            if stripped.startswith("import "):
                imports.append(stripped)
            elif "class " in stripped:
                classes.append(stripped[:80])
            elif stripped.startswith("public ") and "(" in stripped:
                functions.append(stripped[:80])

    summary = f"File: {filename}\nLanguage: {language}\n"
    if imports:
        summary += f"\nImports/Dependencies:\n" + "\n".join(imports[:20])
    if classes:
        summary += f"\nClasses: {', '.join(classes[:15])}"
    if functions:
        summary += f"\nFunctions: {', '.join(functions[:20])}"
    summary += f"\n\nFull code:\n{code}"
    return summary

def parse_code_zip(file_bytes: bytes) -> list[dict]:
    """
    Extract and parse all code files from a ZIP archive.
    Returns list of {filename, language, content, structure}
    """
    files = []
    try:
        with zipfile.ZipFile(io.BytesIO(file_bytes)) as zf:
            for zip_path in zf.namelist():
                # Skip directories and hidden/system files
                parts = zip_path.split("/")
                if any(p in SKIP_DIRS for p in parts):
                    continue
                if os.path.basename(zip_path).startswith("."):
                    continue
                if zip_path.endswith("/"):
                    continue

                filename = os.path.basename(zip_path)
                language = _get_file_language(filename)
                if not language:
                    continue  # skip non-code files

                try:
                    content = zf.read(zip_path).decode("utf-8", errors="ignore")
                    if not content.strip():
                        continue
                    # Limit very large files
                    if len(content) > 50000:
                        content = content[:50000] + "\n... [truncated]"

                    structure = _extract_code_structure(content, language, zip_path)
                    files.append({
                        "zip_path": zip_path,
                        "filename": filename,
                        "language": language,
                        "content":  content,
                        "structure": structure
                    })
                except Exception:
                    continue
    except Exception as e:
        print(f"   ❌ ZIP parsing failed: {e}")

    return files

def parse_single_code_file(file_bytes: bytes, filename: str) -> list[dict]:
    """Parse a single uploaded code file"""
    language = _get_file_language(filename)
    if not language:
        return []
    try:
        content = file_bytes.decode("utf-8", errors="ignore")
        if not content.strip():
            return []
        if len(content) > 50000:
            content = content[:50000] + "\n... [truncated]"
        structure = _extract_code_structure(content, language, filename)
        return [{
            "zip_path": filename,
            "filename": filename,
            "language": language,
            "content":  content,
            "structure": structure
        }]
    except Exception:
        return []

# ── Parsers ────────────────────────────────────────────────
def parse_pdf(file_bytes: bytes) -> tuple[list[str], int]:
    doc   = fitz.open(stream=file_bytes, filetype="pdf")
    pages = [page.get_text() for page in doc]
    return pages, len(pages)

def parse_pptx(file_bytes: bytes, use_vision: bool = True) -> tuple[list[str], int]:
    prs         = Presentation(io.BytesIO(file_bytes))
    slides      = []
    vision_ready = use_vision and _is_vision_available()

    if use_vision and not vision_ready:
        print("⚠️  Vision model not found — text-only extraction.")

    for slide_num, slide in enumerate(prs.slides, start=1):
        text_parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        text_parts.append(t)
            elif hasattr(shape, "text") and shape.text.strip():
                text_parts.append(shape.text.strip())

        text_content = "\n".join(text_parts).strip()
        needs_vision = vision_ready and (
            _slide_has_images(slide) or _slide_text_is_thin(slide)
        )

        if needs_vision:
            print(f"   🔍 Vision on slide {slide_num}...")
            vision_desc = _describe_slide_image(_slide_to_pil_image(slide), slide_num)
            combined = (
                f"{text_content}\n\n[Visual content slide {slide_num}]: {vision_desc}"
                if text_content else
                f"[Slide {slide_num} visual]: {vision_desc}"
            )
            slides.append(combined)
        else:
            slides.append(text_content or f"[Slide {slide_num} — no content]")

    return slides, len(slides)

# ── Main ingestion — PPT/PDF ───────────────────────────────
def ingest_document(team_name: str, filename: str,
                    file_bytes: bytes) -> dict:
    """Ingest PPT or PDF submission"""
    ext          = filename.lower()
    file_type    = "pdf" if ext.endswith(".pdf") else "pptx"
    file_size_kb = len(file_bytes) // 1024

    add_team(team_name)
    doc_id = register_document(team_name, filename, file_type, file_size_kb)

    try:
        update_document_status(doc_id, "processing")
        print(f"\n📄 Parsing {filename} for '{team_name}'...")

        if file_type == "pdf":
            pages, total_pages = parse_pdf(file_bytes)
        else:
            pages, total_pages = parse_pptx(file_bytes, use_vision=True)

        splitter = RecursiveCharacterTextSplitter(
            chunk_size=500, chunk_overlap=50
        )
        all_chunks, all_metadata, all_ids = [], [], []

        for page_num, page_text in enumerate(pages, start=1):
            if not page_text.strip():
                continue
            for chunk_idx, chunk in enumerate(splitter.split_text(page_text)):
                chunk_id = f"doc{doc_id}_p{page_num}_c{chunk_idx}"
                all_chunks.append(chunk)
                all_metadata.append({
                    "team":      team_name,
                    "doc_id":    str(doc_id),
                    "filename":  filename,
                    "file_type": file_type,
                    "page":      page_num,
                    "content_type": "presentation"
                })
                all_ids.append(chunk_id)

        collection = get_collection()
        if all_chunks:
            for i in range(0, len(all_chunks), 100):
                collection.add(
                    documents=all_chunks[i:i+100],
                    metadatas=all_metadata[i:i+100],
                    ids=all_ids[i:i+100]
                )

        update_document_status(
            doc_id, "indexed",
            total_pages=total_pages,
            total_chunks=len(all_chunks)
        )
        print(f"   ✅ {len(all_chunks)} chunks indexed")
        return {"success": True, "doc_id": doc_id,
                "total_pages": total_pages, "total_chunks": len(all_chunks)}

    except Exception as e:
        print(f"   ❌ {e}")
        update_document_status(doc_id, "failed", error_message=str(e))
        return {"success": False, "error": str(e)}

# ── Main ingestion — Code ──────────────────────────────────
def ingest_code(team_name: str, filename: str,
                file_bytes: bytes) -> dict:
    """
    Ingest code — supports:
    - ZIP file containing project
    - Single code file (.py, .js, .java etc)
    """
    file_size_kb = len(file_bytes) // 1024
    add_team(team_name)
    doc_id = register_document(team_name, filename, "code", file_size_kb)

    try:
        update_document_status(doc_id, "processing")
        print(f"\n💻 Parsing code: {filename} for '{team_name}'...")

        # Determine if ZIP or single file
        if filename.lower().endswith(".zip"):
            code_files = parse_code_zip(file_bytes)
        else:
            code_files = parse_single_code_file(file_bytes, filename)

        if not code_files:
            raise ValueError("No supported code files found")

        print(f"   📁 Found {len(code_files)} code files")

        # Use smaller chunks for code to preserve function context
        splitter = RecursiveCharacterTextSplitter(
            chunk_size=800,
            chunk_overlap=100
        )

        all_chunks, all_metadata, all_ids = [], [], []

        for file_info in code_files:
            chunks = splitter.split_text(file_info["structure"])
            for chunk_idx, chunk in enumerate(chunks):
                chunk_id = f"code{doc_id}_{file_info['filename']}_{chunk_idx}"
                # Make chunk_id safe
                chunk_id = re.sub(r"[^a-zA-Z0-9_-]", "_", chunk_id)
                all_chunks.append(chunk)
                all_metadata.append({
                    "team":         team_name,
                    "doc_id":       str(doc_id),
                    "filename":     file_info["zip_path"],
                    "file_type":    "code",
                    "language":     file_info["language"],
                    "page":         0,
                    "content_type": "code"
                })
                all_ids.append(chunk_id)

        collection = get_collection()
        if all_chunks:
            for i in range(0, len(all_chunks), 100):
                collection.add(
                    documents=all_chunks[i:i+100],
                    metadatas=all_metadata[i:i+100],
                    ids=all_ids[i:i+100]
                )

        update_document_status(
            doc_id, "indexed",
            total_pages=len(code_files),
            total_chunks=len(all_chunks)
        )
        print(f"   ✅ {len(all_chunks)} code chunks indexed")
        return {
            "success":      True,
            "doc_id":       doc_id,
            "total_files":  len(code_files),
            "total_chunks": len(all_chunks),
            "languages":    list(set(f["language"] for f in code_files))
        }

    except Exception as e:
        print(f"   ❌ {e}")
        update_document_status(doc_id, "failed", error_message=str(e))
        return {"success": False, "error": str(e)}

# ── Retrieval ──────────────────────────────────────────────
def retrieve_context(query: str, team_name: str = None,
                     n_results: int = 5,
                     content_type: str = None) -> list[dict]:
    """
    Hybrid search: Vector + BM25
    content_type: "presentation", "code", or None (both)
    """
    collection = get_collection()

    # Build where filter
    where = {}
    if team_name:
        where["team"] = team_name
    if content_type:
        where["content_type"] = content_type

    try:
        results = collection.query(
            query_texts=[query],
            n_results=n_results,
            where=where if where else None,
            include=["documents", "metadatas", "distances"]
        )
    except Exception:
        return []

    chunks = []
    if results and results["documents"]:
        for doc, meta, dist in zip(
            results["documents"][0],
            results["metadatas"][0],
            results["distances"][0]
        ):
            chunks.append({
                "text":         doc,
                "filename":     meta.get("filename", ""),
                "page":         meta.get("page", 0),
                "team":         meta.get("team", ""),
                "language":     meta.get("language", ""),
                "content_type": meta.get("content_type", ""),
                "score":        round(1 - dist, 3)
            })

    # BM25 re-rank
    if chunks:
        try:
            from rank_bm25 import BM25Okapi
            tokenised   = [c["text"].lower().split() for c in chunks]
            bm25        = BM25Okapi(tokenised)
            bm25_scores = bm25.get_scores(query.lower().split())
            for i, c in enumerate(chunks):
                c["score"] = round(c["score"] * 0.6 + bm25_scores[i] * 0.4, 3)
            chunks.sort(key=lambda x: x["score"], reverse=True)
        except Exception:
            pass

    return chunks

def delete_team_documents(team_name: str):
    """Remove all ChromaDB entries for a team"""
    collection = get_collection()
    collection.delete(where={"team": team_name})